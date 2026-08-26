#!/usr/bin/env python3
"""Build the OpenAI Jalapeño system-insight deck (16:9)."""

from __future__ import annotations

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.util import Inches, Pt

OUT = "docs/jalapeno/OpenAI_Jalapeno_系统洞察.pptx"
FONT = "微软雅黑"

# Palette
BG = RGBColor(0x0B, 0x10, 0x16)
SURFACE = RGBColor(0x14, 0x1C, 0x26)
CARD = RGBColor(0x1B, 0x25, 0x32)
CARD2 = RGBColor(0x22, 0x2E, 0x3D)
LINE = RGBColor(0x2C, 0x3A, 0x4C)
GREEN = RGBColor(0x3D, 0xDC, 0x84)
GREEN_DIM = RGBColor(0x1A, 0x6B, 0x45)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
BLUE = RGBColor(0x5B, 0xB4, 0xFF)
RED = RGBColor(0xFF, 0x6B, 0x6B)
WHITE = RGBColor(0xF5, 0xF7, 0xFA)
MUTED = RGBColor(0x9A, 0xA7, 0xB5)
DIM = RGBColor(0x6A, 0x78, 0x88)
BLACK = RGBColor(0x08, 0x0C, 0x11)

W = Inches(13.333)
H = Inches(7.5)


def _set_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _line(shape, color, pt=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def _ea_font(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    for tag, attr in (("a:latin", name), ("a:ea", name), ("a:cs", name)):
        node = rPr.find(qn(tag))
        if node is None:
            node = etree.SubElement(rPr, qn(tag))
        node.set("typeface", name)


def _style(run, size, bold=False, color=WHITE, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    run.font.italic = False
    _ea_font(run, name)


def add_textbox(slide, l, t, w, h, text, size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    _style(run, size, bold, color)
    return box


def add_paras(slide, l, t, w, h, items, size=13, color=WHITE, spacing=8,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, bold=False):
    """items: list of str, or (text, size, bold, color)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))
    except Exception:
        pass
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing)
        p.space_before = Pt(0)
        if isinstance(item, tuple):
            text, sz, bd, col = item
        else:
            text, sz, bd, col = item, size, bold, color
        run = p.add_run()
        run.text = text
        _style(run, sz, bd, col)
    return box


def rect(slide, l, t, w, h, fill, line=None, line_pt=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _set_solid(sh, fill)
    if line:
        _line(sh, line, line_pt)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def round_rect(slide, l, t, w, h, fill, line=None, adj=0.08, line_pt=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _set_solid(sh, fill)
    try:
        sh.adjustments[0] = adj
    except Exception:
        pass
    if line:
        _line(sh, line, line_pt)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def accent_bar(slide, l, t, w, h, color=GREEN):
    return rect(slide, l, t, w, h, color)


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, BG)
    return slide


def footer(slide, page, total):
    rect(slide, 0, Inches(7.28), W, Inches(0.22), RGBColor(0x08, 0x0B, 0x10))
    accent_bar(slide, 0, Inches(7.28), Inches(0.18), Inches(0.22), GREEN)
    add_textbox(slide, Inches(0.4), Inches(7.28), Inches(10.5), Inches(0.22),
                "OpenAI Jalapeño 系统洞察  ·  Hot Chips 2026  ·  内部学习资料",
                9, False, DIM, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(11.4), Inches(7.28), Inches(1.6), Inches(0.22),
                f"{page:02d}  /  {total:02d}",
                9, False, MUTED, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)


def title_block(slide, kicker, title, subtitle=None):
    accent_bar(slide, 0, 0, W, Inches(0.06), GREEN)
    add_textbox(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.28),
                kicker, 11, True, GREEN)
    add_textbox(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.48),
                title, 26, True, WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.90), Inches(12.3), Inches(0.32),
                    subtitle, 13, False, MUTED)
        return Inches(1.28)
    return Inches(1.02)


def card(slide, l, t, w, h, fill=CARD, line=LINE):
    return round_rect(slide, l, t, w, h, fill, line, 0.06, 1.0)


def kpi_card(slide, l, t, w, h, value, label, sub=None, value_color=GREEN):
    card(slide, l, t, w, h)
    accent_bar(slide, l, t, Inches(0.08), h, GREEN)
    add_textbox(slide, l + Inches(0.22), t + Inches(0.16), w - Inches(0.32), Inches(0.46),
                value, 22, True, value_color, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_textbox(slide, l + Inches(0.22), t + Inches(0.62), w - Inches(0.32), Inches(0.32),
                label, 13, True, WHITE)
    if sub:
        add_textbox(slide, l + Inches(0.22), t + Inches(0.98), w - Inches(0.32), Inches(1.30),
                    sub, 12, False, MUTED)


def pill(slide, l, t, w, h, text, fill=GREEN, text_color=BLACK, size=11):
    round_rect(slide, l, t, w, h, fill, None, 0.5)
    add_textbox(slide, l, t, w, h, text, size, True, text_color,
                PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def arrow_right(slide, l, t, w, h, fill=GREEN):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    _set_solid(sh, fill)
    sh.shadow.inherit = False
    return sh


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

TOTAL = 24


def s01_cover(prs):
    slide = blank(prs)
    rect(slide, 0, 0, Inches(0.18), H, GREEN)
    rect(slide, 0, Inches(6.55), W, Inches(0.95), RGBColor(0x08, 0x0C, 0x11))
    accent_bar(slide, 0, Inches(6.55), W, Inches(0.06), GREEN)

    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.32),
                "HOT CHIPS 2026  ·  STANFORD  ·  2026.08.25", 13, True, GREEN)
    add_textbox(slide, Inches(0.7), Inches(1.15), Inches(12), Inches(0.9),
                "Jalapeño", 54, True, WHITE)
    add_textbox(slide, Inches(0.7), Inches(2.05), Inches(12), Inches(0.55),
                "OpenAI 推理芯片系统洞察", 28, True, GREEN)
    add_textbox(slide, Inches(0.7), Inches(2.75), Inches(11.5), Inches(0.7),
                "一颗为 LLM Serving 重写的 ASIC，而不是训练 GPU 的缩小版。\n"
                "从问题定义、微架构、机架网络、软件栈到跑分口径，做一次完整拆解。",
                16, False, MUTED)

    # three thesis chips
    items = [
        ("01", "少搬数据", "KV / 权重尽量留本地"),
        ("02", "同质机池", "故意不做 PD 分离"),
        ("03", "电力即吞吐", "先打 tokens / MW"),
    ]
    for i, (n, a, b) in enumerate(items):
        x = Inches(0.7) + i * Inches(4.05)
        card(slide, x, Inches(3.75), Inches(3.85), Inches(1.45))
        add_textbox(slide, x + Inches(0.22), Inches(3.88), Inches(3.4), Inches(0.28),
                    n, 12, True, GREEN)
        add_textbox(slide, x + Inches(0.22), Inches(4.18), Inches(3.4), Inches(0.38),
                    a, 18, True, WHITE)
        add_textbox(slide, x + Inches(0.22), Inches(4.58), Inches(3.4), Inches(0.40),
                    b, 13, False, MUTED)

    add_textbox(slide, Inches(0.7), Inches(6.72), Inches(8), Inches(0.55),
                "OpenAI × Broadcom × Celestica    |    演讲 You Can Just Build Things … Chips",
                13, False, MUTED, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(9.2), Inches(6.72), Inches(3.6), Inches(0.55),
                "内部技术简报  ·  24 页",
                13, False, DIM, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)


def s02_agenda(prs, n):
    slide = blank(prs)
    y0 = title_block(slide, "CONTENTS", "目录", "按「问题 → 架构 → 系统 → 软件 → 成绩 → 战略」读，而不是按参数表读。")
    footer(slide, n, TOTAL)
    chapters = [
        ("01", "定位", "它是什么、为什么现在做、和谁一起做"),
        ("02", "架构哲学", "Prefill/Decode、KV 本地、不做 PD 分离"),
        ("03", "硅与系统", "芯片规格、机架解剖、2048 超节点网络"),
        ("04", "软件栈", "Gluon / Teacup / Codex 写 kernel"),
        ("05", "成绩怎么读", "InferenceX、对标 Blackwell/Rubin、口径陷阱"),
        ("06", "战略含义", "对 OpenAI、Nvidia 与超节点路线的启示"),
    ]
    for i, (num, title, desc) in enumerate(chapters):
        col, row = i % 3, i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = y0 + Inches(0.25) + row * Inches(2.45)
        card(slide, x, y, Inches(4.0), Inches(2.2))
        add_textbox(slide, x + Inches(0.28), y + Inches(0.28), Inches(3.4), Inches(0.4),
                    num, 16, True, GREEN)
        add_textbox(slide, x + Inches(0.28), y + Inches(0.75), Inches(3.4), Inches(0.45),
                    title, 22, True, WHITE)
        add_textbox(slide, x + Inches(0.28), y + Inches(1.28), Inches(3.4), Inches(0.65),
                    desc, 13, False, MUTED)


def s03_thesis(prs, n):
    slide = blank(prs)
    title_block(slide, "THESIS", "先看结论，再看芯片", "Hot Chips 这场的真正内容，不是又一张 FLOPS 海报。")
    footer(slide, n, TOTAL)

    points = [
        ("它解决的是 Serving，不是 Training",
         "Jalapeño 是 OpenAI 第一颗 Intelligence Processor：为在线推理、低延迟 Agent、高并发 ChatGPT 流量设计。训练仍走 Nvidia / AMD / Cerebras。"),
        ("架构选择服从真实流量，而不是服从高峰算力",
         "Prefill 吃算力、Decode 吃带宽、KV 搬家最伤延迟。所以核心不是「更大 systolic」，而是「少搬数据、同质机池、贴近 roofline」。"),
        ("成绩要和口径一起读",
         "公开数字打赢的是 Blackwell，真正的同代对手是 Rubin；而且目前是 STP、无 speculative decoding、无 AgentX。方向对，但还不是终局数字。"),
    ]
    for i, (t, d) in enumerate(points):
        y = Inches(1.35) + i * Inches(1.75)
        card(slide, Inches(0.5), y, Inches(12.35), Inches(1.58))
        round_rect(slide, Inches(0.72), y + Inches(0.48), Inches(0.62), Inches(0.62), GREEN, None, 0.15)
        add_textbox(slide, Inches(0.72), y + Inches(0.48), Inches(0.62), Inches(0.62),
                    f"{i+1:02d}", 16, True, BLACK, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.55), y + Inches(0.22), Inches(10.9), Inches(0.45),
                    t, 18, True, WHITE)
        add_textbox(slide, Inches(1.55), y + Inches(0.72), Inches(10.9), Inches(0.68),
                    d, 14, False, MUTED)


def s04_what_is(prs, n):
    slide = blank(prs)
    title_block(slide, "01  定位", "它是什么，不是什么", "先把边界画清楚，后面的架构选择才读得通。")
    footer(slide, n, TOTAL)

    card(slide, Inches(0.5), Inches(1.35), Inches(6.0), Inches(5.55), RGBColor(0x12, 0x2A, 0x1E))
    accent_bar(slide, Inches(0.5), Inches(1.35), Inches(0.1), Inches(5.55), GREEN)
    add_textbox(slide, Inches(0.85), Inches(1.52), Inches(5.4), Inches(0.4),
                "是", 16, True, GREEN)
    yes = [
        "OpenAI 第一颗自研推理 ASIC，多代平台的 Gen-1",
        "围绕 LLM serving 从零设计：模型、kernel、机架一起做",
        "Broadcom 做硅实现与 Tomahawk 网络，Celestica 做板卡机架",
        "可跑多种模型：GPT-OSS、DeepSeek R1、Kimi K2.5，不是单模型专用芯",
        "自家 captive silicon，不外卖、不出租",
        "2026 年底小规模上量，2027 年才是真正放量",
    ]
    for i, t in enumerate(yes):
        y = Inches(2.05) + i * Inches(0.72)
        round_rect(slide, Inches(0.85), y, Inches(0.28), Inches(0.28), GREEN, None, 0.4)
        add_textbox(slide, Inches(0.85), y, Inches(0.28), Inches(0.28),
                    "✓", 12, True, BLACK, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.28), y - Inches(0.04), Inches(4.95), Inches(0.62),
                    t, 13, False, WHITE)

    card(slide, Inches(6.75), Inches(1.35), Inches(6.1), Inches(5.55), RGBColor(0x2A, 0x16, 0x16))
    accent_bar(slide, Inches(6.75), Inches(1.35), Inches(0.1), Inches(5.55), RED)
    add_textbox(slide, Inches(7.1), Inches(1.52), Inches(5.5), Inches(0.4),
                "不是", 16, True, RED)
    no = [
        "不是训练芯片，也不会替代 OpenAI 的 GPU 训练集群",
        "不是对外售卖的商用加速卡，没有云上实例可租",
        "不是「专门打自家模型」的固化 ASIC（OpenAI 自己强调通用）",
        "不是立刻量产的产品：Hot Chips 上跑的是工程样片 / A0",
        "不是对 Nvidia 的全面宣战，Ho 明确说算力缺口要多方补",
        "不是已经带 speculative decoding / PD 分离的终局系统",
    ]
    for i, t in enumerate(no):
        y = Inches(2.05) + i * Inches(0.72)
        round_rect(slide, Inches(7.1), y, Inches(0.28), Inches(0.28), RED, None, 0.4)
        add_textbox(slide, Inches(7.1), y, Inches(0.28), Inches(0.28),
                    "×", 12, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.53), y - Inches(0.04), Inches(5.05), Inches(0.62),
                    t, 13, False, WHITE)


def s05_why(prs, n):
    slide = blank(prs)
    title_block(slide, "01  定位", "为什么 OpenAI 现在必须自己做芯片", "驱动因素不是「想当芯片公司」，而是 Serving 成本、电力上限和供应结构。")
    footer(slide, n, TOTAL)

    drivers = [
        ("电力是硬约束", "OpenAI 现在卡的是数据中心 MW，不是预算或机房面积。tokens / MW 直接等于收入。Jensen 在 Computex 也说：1 GW 电力下，吞吐每瓦就是营收。"),
        ("推理已经比训练更贵", "ChatGPT / Agent 是持续在线流量。训练是一次性高峰，推理是每天的电费和时延 SLO。垂直整合的最大杠杆在 Serving。"),
        ("GPU 通才的税太重", "通用 GPU 要兼容训练、图形、CUDA 生态。对 decode 带宽、KV 局部性、小 batch 延迟，通才架构付了固定开销税。"),
        ("供应不能只押一家", "OpenAI 仍会买 Nvidia / AMD / Cerebras。自研 ASIC 是把一部分推理从高毛利 GPU 挪到更可控的 Broadcom 供应链上。"),
        ("全栈才能改问题定义", "只有同时拥有模型、kernel、调度和硅，才能决定：KV 放哪、要不要 PD 分离、哪些延迟必须从硬件里删掉。"),
        ("Agent 改变流量形状", "知识 → 推理 → Agent 三代模型，I/O 比、cache hit、投机接受率都在变。固定拆分的异构机池会很快过时。"),
    ]
    for i, (t, d) in enumerate(drivers):
        col, row = i % 3, i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.38) + row * Inches(2.7)
        card(slide, x, y, Inches(4.0), Inches(2.5))
        add_textbox(slide, x + Inches(0.22), y + Inches(0.2), Inches(3.55), Inches(0.28),
                    f"0{i+1}", 12, True, GREEN)
        add_textbox(slide, x + Inches(0.22), y + Inches(0.5), Inches(3.55), Inches(0.45),
                    t, 16, True, WHITE)
        add_textbox(slide, x + Inches(0.22), y + Inches(1.02), Inches(3.55), Inches(1.25),
                    d, 12, False, MUTED)


def s06_timeline(prs, n):
    slide = blank(prs)
    title_block(slide, "01  定位", "时间线与分工", "OpenAI 设计，Broadcom 落地硅和网络，Celestica 做可量产的机架。")
    footer(slide, n, TOTAL)

    events = [
        ("2024 H2", "组队启动", "从招人到定义架构"),
        ("2025.11", "CoWoS 流片", "OpenAI 称设计到 tape-out 约 9 个月"),
        ("2026.06", "正式亮相", "与 Broadcom 发布 Jalapeño"),
        ("2026.08", "Hot Chips", "首次公开规格、功耗与跑分"),
        ("2026 年底", "小规模部署", "工程样片 → 少量上架"),
        ("2027", "放量之年", "大部分产出集中在年底"),
    ]
    # timeline line
    rect(slide, Inches(0.7), Inches(2.22), Inches(12.0), Inches(0.04), GREEN)
    for i, (when, what, note) in enumerate(events):
        x = Inches(0.55) + i * Inches(2.1)
        round_rect(slide, x + Inches(0.78), Inches(2.10), Inches(0.28), Inches(0.28), GREEN, None, 0.5)
        add_textbox(slide, x, Inches(1.45), Inches(1.9), Inches(0.55),
                    when, 12, True, GREEN, PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.50), Inches(1.9), Inches(0.4),
                    what, 14, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.90), Inches(1.9), Inches(0.55),
                    note, 11, False, MUTED, PP_ALIGN.CENTER)

    # partner cards
    partners = [
        ("OpenAI", "架构与软件", "芯片定义、kernel、serving 引擎 Teacup、模型与 Codex 协同设计。演讲：Richard Ho / Ravi Narayanaswami / Chris Leary。"),
        ("Broadcom", "硅实现 + 网络", "N3P 计算 die、N3E I/O die、Tomahawk 6 交换、SerDes。把「能画的架构」变成「能量产的硅」。"),
        ("Celestica", "板卡 / 机架", "Katsu 主机柜、Vindaloo ASIC 托盘、Chana 交换托盘，以及 dock-to-rack 部署。"),
    ]
    for i, (name, role, desc) in enumerate(partners):
        x = Inches(0.5) + i * Inches(4.2)
        card(slide, x, Inches(3.65), Inches(4.0), Inches(3.2))
        pill(slide, x + Inches(0.22), Inches(3.85), Inches(1.7), Inches(0.32), role, GREEN, BLACK, 10)
        add_textbox(slide, x + Inches(0.22), Inches(4.30), Inches(3.55), Inches(0.45),
                    name, 20, True, WHITE)
        add_textbox(slide, x + Inches(0.22), Inches(4.85), Inches(3.55), Inches(1.7),
                    desc, 13, False, MUTED)


def s07_section_arch(prs, n):
    slide = blank(prs)
    rect(slide, 0, 0, Inches(0.18), H, GREEN)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.4),
                "CHAPTER 02", 14, True, GREEN)
    add_textbox(slide, Inches(0.8), Inches(2.65), Inches(11.5), Inches(0.8),
                "架构哲学：先改问题，再做硅", 32, True, WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11.2), Inches(1.1),
                "Jalapeño 真正特别的地方，不是 13.4 PFLOPS 这个数字，\n"
                "而是它拒绝了行业里两件「默认正确」的事：通用 GPU 税，以及 Prefill/Decode 分离。",
                16, False, MUTED)
    footer(slide, n, TOTAL)


def s08_prefill_decode(prs, n):
    slide = blank(prs)
    title_block(slide, "02  架构", "LLM 推理的三个瓶颈，不是一个", "一颗只打高峰算力的芯片，会在另外两个瓶颈上把优势吐回去。")
    footer(slide, n, TOTAL)

    phases = [
        ("Prefill", "算力密集", BLUE,
         "把 prompt 一次性编进 KV。\n大矩阵、高算力利用率。\n长上下文时这里会爆。",
         "优化点：矩阵引擎、小 shape 不掉崖"),
        ("Decode", "带宽密集", GREEN,
         "逐 token 生成。每步都要读权重和 KV。\n算力经常闲着，HBM 才是上限。",
         "优化点：HBM4 带宽、本地 KV、少搬数据"),
        ("通信 / 调度", "延迟税", AMBER,
         "KV 跨芯片、跨机、跨 PD 池搬家。\nlaunch / barrier / 排队都会变成固定开销。",
         "优化点：删固定延迟，让小 batch 也贴 roofline"),
    ]
    for i, (name, tag, col, body, foot) in enumerate(phases):
        x = Inches(0.5) + i * Inches(4.2)
        card(slide, x, Inches(1.4), Inches(4.0), Inches(3.55))
        accent_bar(slide, x, Inches(1.4), Inches(4.0), Inches(0.08), col)
        add_textbox(slide, x + Inches(0.22), Inches(1.6), Inches(3.55), Inches(0.4),
                    name, 20, True, WHITE)
        pill(slide, x + Inches(0.22), Inches(2.08), Inches(1.7), Inches(0.30), tag, col, BLACK, 11)
        add_textbox(slide, x + Inches(0.22), Inches(2.52), Inches(3.55), Inches(1.45),
                    body, 13, False, MUTED)
        add_textbox(slide, x + Inches(0.22), Inches(4.15), Inches(3.55), Inches(0.55),
                    foot, 12, True, col)
        if i < 2:
            arrow_right(slide, x + Inches(3.85), Inches(2.9), Inches(0.42), Inches(0.22), col)

    card(slide, Inches(0.5), Inches(5.15), Inches(12.35), Inches(1.75))
    add_textbox(slide, Inches(0.75), Inches(5.32), Inches(11.9), Inches(0.35),
                "OpenAI 的原话，其实就一句", 13, True, GREEN)
    add_textbox(slide, Inches(0.75), Inches(5.70), Inches(11.9), Inches(0.95),
                "「只擅长其中一个阶段的系统，会在等数据、搬模型状态时把优势输掉。」所以 Jalapeño 的目标不是 Prefill 芯片或 Decode 芯片，"
                "而是一颗能按阶段切换计算 / 内存 / 网络组合、并且把 KV 明确放在本地的加速器。",
                14, False, WHITE)


def s09_locality(prs, n):
    slide = blank(prs)
    title_block(slide, "02  架构", "核心设计：少搬数据，KV 留本地", "内存层次做减法，比把 systolic 做加法更重要。")
    footer(slide, n, TOTAL)

    # left: slice diagram
    card(slide, Inches(0.5), Inches(1.35), Inches(7.4), Inches(5.55))
    add_textbox(slide, Inches(0.75), Inches(1.5), Inches(7.0), Inches(0.35),
                "Slice 结构（示意）", 14, True, GREEN)

    for i in range(3):
        x = Inches(0.85) + i * Inches(2.25)
        round_rect(slide, x, Inches(2.05), Inches(2.05), Inches(2.55), CARD2, GREEN, 0.08, 1.25)
        add_textbox(slide, x, Inches(2.18), Inches(2.05), Inches(0.32),
                    f"Slice {i}", 12, True, GREEN, PP_ALIGN.CENTER)
        round_rect(slide, x + Inches(0.18), Inches(2.58), Inches(1.7), Inches(0.85), SURFACE, LINE)
        add_textbox(slide, x + Inches(0.18), Inches(2.58), Inches(1.7), Inches(0.85),
                    "Core\nOoO + L1", 12, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        round_rect(slide, x + Inches(0.18), Inches(3.55), Inches(1.7), Inches(0.82), RGBColor(0x1A, 0x3A, 0x2A), GREEN)
        add_textbox(slide, x + Inches(0.18), Inches(3.55), Inches(1.7), Inches(0.82),
                    "本地 HBM 切片\n低延迟直视", 12, True, GREEN, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    round_rect(slide, Inches(0.85), Inches(4.80), Inches(6.55), Inches(0.55), RGBColor(0x1A, 0x3A, 0x2A), GREEN)
    add_textbox(slide, Inches(0.85), Inches(4.80), Inches(6.55), Inches(0.55),
                "专用 Collective 网络  ·  slice 间同步（可与计算重叠的 TP 通信）",
                12, True, GREEN, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    round_rect(slide, Inches(0.85), Inches(5.50), Inches(6.55), Inches(0.55), SURFACE, LINE)
    add_textbox(slide, Inches(0.85), Inches(5.50), Inches(6.55), Inches(0.55),
                "通用 NoC  ·  访问 scale-up 网络 / 一般通信",
                12, True, MUTED, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.75), Inches(6.15), Inches(6.9), Inches(0.55),
                "权重和 KV 精心放置后，核间同步被限制在已知的高带宽通信上。",
                12, False, MUTED)

    # right points
    notes = [
        ("对比 GPU", "GPU 访存要穿过复杂内存系统，大延迟必须靠更大 shape / 更大 batch 摊销。Jalapeño 把层次压扁，小 batch 也能靠近峰值。"),
        ("对比 TPU 式大 systolic", "同样是 weight-stationary 脉动阵列，但明确支持更小矩阵维，减少 padding / tiling 掉崖。"),
        ("对比 PD 分离", "Prefill 产出的 KV 立刻被 Decode 需要。搬走它等于把局部性拆掉，再买一次网络和排队。"),
    ]
    for i, (t, d) in enumerate(notes):
        y = Inches(1.35) + i * Inches(1.85)
        card(slide, Inches(8.1), y, Inches(4.75), Inches(1.72))
        add_textbox(slide, Inches(8.32), y + Inches(0.15), Inches(4.35), Inches(0.35),
                    t, 14, True, GREEN)
        add_textbox(slide, Inches(8.32), y + Inches(0.52), Inches(4.35), Inches(1.05),
                    d, 12, False, MUTED)


def s10_microarch(prs, n):
    slide = blank(prs)
    title_block(slide, "02  架构", "微架构：删固定延迟，让小 batch 也贴顶", "和多数加速器相反，它用乱序核 + L1，而不是纯软件 scratchpad。")
    footer(slide, n, TOTAL)

    items = [
        ("计算", "MXFP 数值格式 + weight-stationary 脉动阵列（像 TPU）。但支持更小 shape，避免大阵列在奇怪维数上的性能悬崖。"),
        ("标量 / 向量", "64-bit 标量核 + FP32/INT32 向量核。不是纯矩阵黑盒，保留可编程性，所以能跑 Doom，也能快速搬新模型。"),
        ("控制", "乱序核 + L1 cache。刻意避开 barrier / 软件 scratchpad 的固定开销。代价是更依赖预取质量。"),
        ("冗余", "托盘级冗余，核级 / channel 级 yield harvesting。第一代 ASIC 能上量，这块和峰值 FLOPS 一样关键。"),
        ("I/O die", "N3E I/O chiplet：32×800G SerDes。24 lane 本地 scale-up（600 GB/s），8 lane 全局（200 GB/s）。主机侧 PCIe Gen5。"),
        ("AI 辅助设计", "OpenAI 称 AI 把 SIMD 面积压了 8%、矩阵引擎面积压了 10%，并改善时序和功耗。B0 相对 A0 还有约 25% perf/W。"),
    ]
    for i, (t, d) in enumerate(items):
        col, row = i % 3, i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.35) + row * Inches(2.7)
        card(slide, x, y, Inches(4.0), Inches(2.5))
        add_textbox(slide, x + Inches(0.22), y + Inches(0.22), Inches(3.55), Inches(0.4),
                    t, 16, True, GREEN)
        add_textbox(slide, x + Inches(0.22), y + Inches(0.72), Inches(3.55), Inches(1.55),
                    d, 13, False, WHITE)


def s11_no_pdd(prs, n):
    slide = blank(prs)
    title_block(slide, "02  架构", "最反直觉的决定：不做 Prefill / Decode 分离", "对 GPU 几乎是标配的 PDD，在 Jalapeño 上被明确拒绝。")
    footer(slide, n, TOTAL)

    # two columns
    card(slide, Inches(0.5), Inches(1.35), Inches(6.05), Inches(3.55))
    add_textbox(slide, Inches(0.75), Inches(1.5), Inches(5.6), Inches(0.35),
                "行业默认：PD 分离", 16, True, AMBER)
    for i, t in enumerate([
        "Prefill 池吃算力，Decode 池吃带宽，各自调到最优点",
        "流量 I/O 比一变，就会一边排队、一边空转",
        "KV 必须跨池传输：多一次带宽、同步、排队和故障域",
        "Draft / Verify 若再拆开，投机解码会变成分布式协议",
    ]):
        add_textbox(slide, Inches(0.75), Inches(2.0) + i * Inches(0.65), Inches(5.55), Inches(0.6),
                    f"·  {t}", 13, False, MUTED)

    card(slide, Inches(6.75), Inches(1.35), Inches(6.1), Inches(3.55), RGBColor(0x12, 0x2A, 0x1E))
    add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.35),
                "Jalapeño：同质可互换机池", 16, True, GREEN)
    for i, t in enumerate([
        "每个设备都能接下一个请求，不存在「分错池」的整卡空闲",
        "KV 留在产生它的地方，局部性变成延迟和功耗优势",
        "Draft 和主模型走同一低延迟 fabric，投机解码才值得做",
        "Agent 流量的 I/O 比会持续漂移，同质池是在买期权",
    ]):
        add_textbox(slide, Inches(7.0), Inches(2.0) + i * Inches(0.65), Inches(5.6), Inches(0.6),
                    f"·  {t}", 13, False, WHITE)

    card(slide, Inches(0.5), Inches(5.08), Inches(12.35), Inches(1.82))
    add_textbox(slide, Inches(0.75), Inches(5.22), Inches(11.9), Inches(0.32),
                "对超节点 / PD 分离路线的直接含义", 14, True, GREEN)
    add_textbox(slide, Inches(0.75), Inches(5.58), Inches(11.9), Inches(1.1),
                "PD 分离不是免费午餐：它在「流量稳定、需要很大 phase-specific batch 才能喂饱 GPU」时很划算；"
                "一旦芯片本身就能在小 batch / 全曲线贴近 roofline，再拆池反而损失局部性。"
                "超节点该优先保证的是 KV 与 TP/EP 的高频域，而不是先把 P 和 D 在拓扑上切开。",
                14, False, WHITE)


def s12_section_sys(prs, n):
    slide = blank(prs)
    rect(slide, 0, 0, Inches(0.18), H, GREEN)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.4),
                "CHAPTER 03", 14, True, GREEN)
    add_textbox(slide, Inches(0.8), Inches(2.65), Inches(11.5), Inches(0.8),
                "从一颗硅，到 2048 卡超节点", 32, True, WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11.2), Inches(1.1),
                "Jalapeño 的系统形态并不意外：rack-scale，对标 NVL72 / Helios。\n"
                "真正要看的是：单卡带宽、机架功耗，以及 16 柜一套的全局 scale-up。",
                16, False, MUTED)
    footer(slide, n, TOTAL)


def s13_chip_specs(prs, n):
    slide = blank(prs)
    title_block(slide, "03  硅与系统", "单芯片规格（B0 / 对外口径）", "TSMC N3P 整幅计算 die + N3E I/O die + HBM4。TDP 刻意不高。")
    footer(slide, n, TOTAL)

    specs = [
        ("13.4 PFLOPS", "MXFP4 峰值算力", "同节点 Rubin 计算 die 约 17.5 PF dense NVFP4，但 TDP 高得多"),
        ("216 GB", "HBM4 容量", "约 6 堆 12-Hi。早期采用 HBM4，pin 速约 10 Gbps"),
        ("15.4 TB/s", "封装带宽", "高于在役 HBM3E 加速器；SA 认为可能略高于 Rubin 的 9.6 Gbps"),
        ("700 W", "TDP / 持续更低", "Hot Chips：测试负载持续功耗 ≤ 550 W"),
        ("N3P + N3E", "工艺拆分", "计算 die 走 N3P；I/O chiplet 走 N3E，32×800G SerDes"),
        ("PCIe Gen5", "主机连接", "连 x86 host。scale-up 不走 PCIe，走专用 SerDes + Tomahawk"),
    ]
    for i, (v, l, s) in enumerate(specs):
        col, row = i % 3, i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.35) + row * Inches(2.7)
        kpi_card(slide, x, y, Inches(4.0), Inches(2.5), v, l, s)


def s14_rack(prs, n):
    slide = blank(prs)
    title_block(slide, "03  硅与系统", "机架与 Pod：数字要一起看", "单卡好看不够。Serving 的真实单位已经是机架，甚至是 16 柜一个域。")
    footer(slide, n, TOTAL)

    kpis = [
        ("128", "卡 / 机架", "1.7 EFLOPS  4-bit\n27.5 TB HBM4  ·  ~2 PB/s"),
        ("2,048", "卡 / Pod", "16 机架组成一个全局 scale-up 域"),
        ("~160 kW", "双柜系统功耗", "Host 柜约 31 kW 生产功耗\nASIC 柜约 130 kW"),
        ("×16", "全局扩展", "Local 128 all-to-all\nGlobal 2048，铜 + 光 + OCS"),
    ]
    for i, (v, l, s) in enumerate(kpis):
        x = Inches(0.5) + i * Inches(3.2)
        card(slide, x, Inches(1.35), Inches(3.05), Inches(2.55))
        add_textbox(slide, x + Inches(0.15), Inches(1.5), Inches(2.75), Inches(0.7),
                    v, 28, True, GREEN, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(2.2), Inches(2.75), Inches(0.4),
                    l, 14, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(2.6), Inches(2.75), Inches(1.05),
                    s, 12, False, MUTED, PP_ALIGN.CENTER)

    # comparison
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.35),
                "和当代 GPU 机架怎么比（公开口径，需打折）", 14, True, WHITE)

    rows = [
        ("", "Jalapeño 机架", "GB200 / GB300 NVL72", "Helios / Rubin 方向"),
        ("定位", "纯推理", "训练 + 推理", "训练 + 推理"),
        ("算力", "1.7 EF  4-bit", "同代更高峰值", "峰值更高，TDP 也更高"),
        ("内存带宽", "约 2 PB/s", "公开对比约 85% 于 Jalapeño", "HBM4，同代竞品"),
        ("功耗哲学", "700W，不追峰值 TDP", "高 TDP 换高峰 FLOPS", "Max-Q 也到 1800W 量级"),
    ]
    col_w = [Inches(2.0), Inches(3.4), Inches(3.5), Inches(3.45)]
    for r, row in enumerate(rows):
        y = Inches(4.5) + r * Inches(0.42)
        x = Inches(0.5)
        for c, cell in enumerate(row):
            fill = RGBColor(0x1A, 0x3A, 0x2A) if r == 0 else (CARD if r % 2 else SURFACE)
            color = GREEN if r == 0 or c == 1 else WHITE
            rect(slide, x, y, col_w[c], Inches(0.42), fill, LINE, 0.75)
            add_textbox(slide, x + Inches(0.08), y, col_w[c] - Inches(0.1), Inches(0.42),
                        cell, 11, r == 0 or c == 0, color, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            x += col_w[c]


def s15_anatomy(prs, n):
    slide = blank(prs)
    title_block(slide, "03  硅与系统", "机架解剖：Katsu / Vindaloo / Chana", "整柜菜谱式命名。系统级设计与 Celestica 合作，scale-up 铜缆背板类似 Nvidia Oberon。")
    footer(slide, n, TOTAL)

    blocks = [
        ("Katsu  主机柜", GREEN,
         "16 个主机托盘，一对一对应 16 个 ASIC 托盘。",
         ["每托盘 2× AMD EPYC Turin", "1.5 TB DRAM / 托盘", "2× E1.S + 2× M.2", "前端网络 400G（2×200G）", "8 根外部 PCIe DAC 横连 Vindaloo"]),
        ("Vindaloo  ASIC 柜", AMBER,
         "16 个 ASIC 托盘，每盘 8 颗 Jalapeño，整柜 128 卡。",
         ["每盘 8 ASIC", "铜缆背板连 Chana", "推理主计算面", "生产功耗大头在这一柜", "与 Host 柜组成双柜系统"]),
        ("Chana  交换托盘", BLUE,
         "8 个 scale-up 交换托盘：6 本地 + 2 全局。",
         ["本地：6× Tomahawk 6 102.4T", "全局：每盘或为 2× TH6", "Local 域 128 卡 all-to-all", "Global 域最多 16 柜 / 2048 卡", "全局走 1.6T 光模块 + OCS"]),
    ]
    for i, (title, col, lead, bullets) in enumerate(blocks):
        x = Inches(0.5) + i * Inches(4.2)
        card(slide, x, Inches(1.35), Inches(4.0), Inches(5.55))
        accent_bar(slide, x, Inches(1.35), Inches(4.0), Inches(0.08), col)
        add_textbox(slide, x + Inches(0.22), Inches(1.55), Inches(3.55), Inches(0.45),
                    title, 16, True, WHITE)
        add_textbox(slide, x + Inches(0.22), Inches(2.05), Inches(3.55), Inches(0.7),
                    lead, 13, False, MUTED)
        for j, b in enumerate(bullets):
            add_textbox(slide, x + Inches(0.22), Inches(2.85) + j * Inches(0.55), Inches(3.55), Inches(0.5),
                        f"·  {b}", 13, False, WHITE)


def s16_network(prs, n):
    slide = blank(prs)
    title_block(slide, "03  硅与系统", "Scale-up：128 本地 + 2048 全局", "网络只占系统成本约 10%，却买到了未来 10–20T 参数 / 百万级上下文的期权。")
    footer(slide, n, TOTAL)

    # two domain cards
    card(slide, Inches(0.5), Inches(1.35), Inches(6.05), Inches(3.7))
    pill(slide, Inches(0.75), Inches(1.55), Inches(1.5), Inches(0.3), "LOCAL", GREEN, BLACK, 11)
    add_textbox(slide, Inches(0.75), Inches(1.95), Inches(5.55), Inches(0.4),
                "机架内 128 卡全互联", 18, True, WHITE)
    stats = [("4.8 Tb/s", "每卡单向本地带宽"), ("6 × 102.4T", "Tomahawk 6 本地交换"), ("6,144 DP", "每柜本地无源铜缆对数")]
    for i, (a, b) in enumerate(stats):
        y = Inches(2.5) + i * Inches(0.7)
        add_textbox(slide, Inches(0.75), y, Inches(2.2), Inches(0.55), a, 16, True, GREEN, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.0), y, Inches(3.2), Inches(0.55), b, 13, False, MUTED, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    card(slide, Inches(6.75), Inches(1.35), Inches(6.1), Inches(3.7))
    pill(slide, Inches(7.0), Inches(1.55), Inches(1.6), Inches(0.3), "GLOBAL", BLUE, BLACK, 11)
    add_textbox(slide, Inches(7.0), Inches(1.95), Inches(5.55), Inches(0.4),
                "16 柜 / 2048 卡一个域", 18, True, WHITE)
    stats2 = [("1.6 Tb/s", "每卡单向全局带宽"), ("8-rail", "rail-only + 机架内 OCS"), ("铜 + 1.6T 光", "背板出柜再进光开关")]
    for i, (a, b) in enumerate(stats2):
        y = Inches(2.5) + i * Inches(0.7)
        add_textbox(slide, Inches(7.0), y, Inches(2.2), Inches(0.55), a, 16, True, BLUE, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(9.25), y, Inches(3.3), Inches(0.55), b, 13, False, MUTED, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    card(slide, Inches(0.5), Inches(5.2), Inches(12.35), Inches(1.7))
    add_textbox(slide, Inches(0.75), Inches(5.35), Inches(11.9), Inches(0.32),
                "读法", 13, True, GREEN)
    add_textbox(slide, Inches(0.75), Inches(5.72), Inches(11.9), Inches(0.95),
                "Local 域负责机架内 TP/EP 高频通信；Global 域把 16 柜收成一个 2048 卡的 scale-up 世界。"
                "这和「超节点优先容纳高频并行组」是同一类判断：能留在铜背板里的，就不要先打到 scale-out 以太网。"
                "OpenAI 还在和 neocloud 收集可靠性数据，瓶颈将从芯片转向部署、监控和弹性。",
                13, False, WHITE)


def s17_software(prs, n):
    slide = blank(prs)
    title_block(slide, "04  软件", "软件才是这颗芯片能打的原因", "SemiAnalysis 的判断：硬件未必强过 Rubin，但软件 bring-up 明显更快。")
    footer(slide, n, TOTAL)

    layers = [
        ("Gluon", "编程模型", "基于 Triton 的 kernel 语言，保留 SPMD，但暴露更底层抽象。核心是 Linear Layouts：用代数描述「硬件资源 ↔ tensor 元素」映射，可证明的 layout 变换和最优 swizzle。"),
        ("Teacup", "Serving 引擎", "内部推理引擎。Persistent thread：程序员而不是硬件调度器分配 tile。配合 TensorInfo / 预取 / 解耦乱序单元。"),
        ("Gigakernel", "减少 launch 税", "单 mega-kernel 在设备上循环，砍 CPU 开销和 launch 延迟。方向是把运行时搬进芯片。"),
        ("Codex", "写 kernel 的人", "早期人在环，随后用内部加强版 Codex。kernel 可到约 3000 行汇编级代码，带正确性检查和自定义 sanitizer。InferenceX 上的 MLA kernel 甚至是 Codex 在无人介入下写出的。"),
    ]
    for i, (name, role, desc) in enumerate(layers):
        y = Inches(1.32) + i * Inches(1.38)
        card(slide, Inches(0.5), y, Inches(12.35), Inches(1.26))
        round_rect(slide, Inches(0.72), y + Inches(0.32), Inches(2.15), Inches(0.62), GREEN, None, 0.12)
        add_textbox(slide, Inches(0.72), y + Inches(0.32), Inches(2.15), Inches(0.62),
                    name, 14, True, BLACK, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.1), y + Inches(0.12), Inches(9.4), Inches(0.35),
                    role, 13, True, GREEN)
        add_textbox(slide, Inches(3.1), y + Inches(0.48), Inches(9.4), Inches(0.68),
                    desc, 13, False, WHITE)


def s18_ai_design(prs, n):
    slide = blank(prs)
    title_block(slide, "04  软件", "用模型做芯片，再用芯片跑模型", "速度本身就是架构声明：从零软件栈起步，三个月硅上 bring-up 就能打公开基准。")
    footer(slide, n, TOTAL)

    kpis = [
        ("~9 个月", "设计到 tape-out", "OpenAI 自称高性能 ASIC 最快周期"),
        ("~3 个月", "A0 硅上 bring-up", "从零软件栈跑通 InferenceX"),
        ("8 天", "TP8 → TP32", "从单系统扩到整柜大规模模型"),
        ("< 2 周", "吞吐翻倍以上", "同一交互性点位，kernel 迭代极快"),
    ]
    for i, (v, l, s) in enumerate(kpis):
        x = Inches(0.5) + i * Inches(3.2)
        card(slide, x, Inches(1.35), Inches(3.05), Inches(2.15))
        add_textbox(slide, x + Inches(0.12), Inches(1.5), Inches(2.8), Inches(0.55),
                    v, 22, True, GREEN, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.12), Inches(2.05), Inches(2.8), Inches(0.4),
                    l, 13, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.12), Inches(2.45), Inches(2.8), Inches(0.8),
                    s, 12, False, MUTED, PP_ALIGN.CENTER)

    demos = [
        ("chilisim", "仿真器与硅误差约 5%，用固定位宽 trace bus。先仿真再上硅。"),
        ("Raiku / Spark", "内部模型 5.3 Codex Spark，演示 1.2 ms TPOT。"),
        ("Doom 36 FPS", "Codex 把 Doom 迁到 Jalapeño，用来证明可编程性。"),
        ("B0 已在 fab", "相对 A0 约 +25% perf/W。对外 13.4 PF 数字是 B0。"),
    ]
    for i, (t, d) in enumerate(demos):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(3.7) + row * Inches(1.4)
        card(slide, x, y, Inches(6.2), Inches(1.25))
        add_textbox(slide, x + Inches(0.22), y + Inches(0.15), Inches(5.75), Inches(0.35),
                    t, 14, True, GREEN)
        add_textbox(slide, x + Inches(0.22), y + Inches(0.52), Inches(5.75), Inches(0.55),
                    d, 13, False, WHITE)


def s19_bench(prs, n):
    slide = blank(prs)
    title_block(slide, "05  成绩", "公开跑分：同一套架构既要吞吐也要延迟", "Benchmark：SemiAnalysis InferenceX。模型：GPT-OSS-120B / DeepSeek R1 / Kimi K2.5。")
    footer(slide, n, TOTAL)

    nums = [
        ("1.5–1.9×", "峰值吞吐", "更多 AI work / 单位功耗"),
        ("1.7–3.6×", "端到端延迟", "同模型更低 e2e latency"),
        ("2.1–4.1×", "超低延迟档", "交互式 / 低并发更明显"),
        (">700", "tok/s/user", "DeepSeek R1，concurrency = 1"),
    ]
    for i, (v, l, s) in enumerate(nums):
        x = Inches(0.5) + i * Inches(3.2)
        card(slide, x, Inches(1.35), Inches(3.05), Inches(2.2))
        add_textbox(slide, x + Inches(0.1), Inches(1.5), Inches(2.85), Inches(0.7),
                    v, 24, True, GREEN, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(2.2), Inches(2.85), Inches(0.4),
                    l, 14, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(2.6), Inches(2.85), Inches(0.7),
                    s, 12, False, MUTED, PP_ALIGN.CENTER)

    bullets = [
        ("没有 MTP / speculative decoding", "Jalapeño 用单 token 预测；对手常用 MTP。这让对比更「干净」，也让 Jalapeño 的数字偏保守。"),
        ("没有 PD 分离", "对手 GPU 机架常靠 PDD 抬吞吐。Jalapeño 用同质池打赢，说明局部性和低固定延迟确实在起作用。"),
        ("跨模型，不只是自家模型", "R1 / Kimi K2.5 / GPT-OSS 都能跑，GSM8k 与 Nvidia 持平。Ho：这是通用、灵活的加速器。"),
        ("软件还在爬坡", "两周内吞吐翻倍、8 天打通 TP32。现在看到的是下限，不是上限。"),
    ]
    for i, (t, d) in enumerate(bullets):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(3.75) + row * Inches(1.4)
        card(slide, x, y, Inches(6.2), Inches(1.25))
        add_textbox(slide, x + Inches(0.22), y + Inches(0.14), Inches(5.75), Inches(0.32),
                    t, 14, True, GREEN)
        add_textbox(slide, x + Inches(0.22), y + Inches(0.48), Inches(5.75), Inches(0.62),
                    d, 12, False, WHITE)


def s20_caveats(prs, n):
    slide = blank(prs)
    title_block(slide, "05  成绩", "这些数字该怎么读", "方向很强，但把 Jalapeño 写成「已经全面超越 Rubin」还太早。")
    footer(slide, n, TOTAL)

    rows = [
        ("对标对象", "公开对比主要打 Blackwell GB200/GB300 NVL72。同代真正对手是 Rubin（HBM4，已开始给客户出货）。SA 认为 Jalapeño 在 tok/s/MW 上仍能打过 Rubin 的公开 MTP 数字，但 Rubin 软件同样未成熟。"),
        ("工作负载", "当前是 8k1k 单轮。没有 AgentX：多轮、长上下文、prefix cache、路由、KV offload 这些生产痛点还没被公开压测。"),
        ("模型体量", "GPT-OSS / R1 / Kimi K2.5 不小，但不是开放前沿最大模型（如 DeepSeek V4 Pro、Kimi K3）。新大模型在新芯片上的搬迁成本被低估。"),
        ("硅版本", "公开成绩来自 A0；B0 在 fab，宣称 +25% perf/W。量产爬坡在 2027，大部分产出年底。工程样片 ≠ 数据中心里的稳定机群。"),
        ("数字来源", "数字由 OpenAI 提供。SemiAnalysis 在实验室现场核对了部分 InferenceX，但未跑全套，也未见 AgentX。"),
        ("TCO", "tokens / $ 与 Rubin 接近；Jalapeño 还没上投机解码（SA 估 MTP 可再降 3–5× 成本）。一部分优势来自 Broadcom 毛利低于 Nvidia，不只是微架构。"),
    ]
    for i, (t, d) in enumerate(rows):
        y = Inches(1.26) + i * Inches(0.92)
        card(slide, Inches(0.5), y, Inches(12.35), Inches(0.86))
        add_textbox(slide, Inches(0.7), y, Inches(2.05), Inches(0.86),
                    t, 13, True, AMBER, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.85), y + Inches(0.08), Inches(9.75), Inches(0.70),
                    d, 12, False, WHITE, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def s21_strategy(prs, n):
    slide = blank(prs)
    title_block(slide, "06  战略", "这意味着什么", "对 OpenAI 是补推理产能；对行业是「CUDA 护城河还在不在」的压力测试。")
    footer(slide, n, TOTAL)

    items = [
        ("对 OpenAI",
         "把产品、模型、芯片、内存做成多代共设计。Gen-2 已深入开发，Gen-3 已启动。下一步目标是 100 MW 级部署——瓶颈转向制造、上架、监控、弹性。"),
        ("对 Nvidia / AMD",
         "训练盘仍然稳。推理盘出现一个「不卖卡、只服务自己流量」的顶级客户自研芯。CUDA 的护城河若只靠生态惯性，会被「Codex 写 kernel + 干净架构」打薄。"),
        ("对 ASIC 玩家",
         "Meta / Microsoft 做了更久却没做成，说明便宜不是充分条件。OpenAI 的差异是：有真实 serving 负载、有 kernel 团队、有模型写代码。"),
        ("对超节点路线",
         "可迁移的三条：① 高频域做大（128 / 2048）；② 优先保 KV 局部性，而不是先 PD 切开；③ 用软件迭代速度而不是一次把 ISA 做完美。"),
    ]
    for i, (t, d) in enumerate(items):
        y = Inches(1.32) + i * Inches(1.38)
        card(slide, Inches(0.5), y, Inches(12.35), Inches(1.26))
        add_textbox(slide, Inches(0.75), y + Inches(0.16), Inches(11.9), Inches(0.32),
                    t, 15, True, GREEN)
        add_textbox(slide, Inches(0.75), y + Inches(0.52), Inches(11.9), Inches(0.6),
                    d, 13, False, WHITE)


def s22_risks(prs, n):
    slide = blank(prs)
    title_block(slide, "06  战略", "风险与未解问题", "第一代芯片好看，不自动等于 2027 年的机群好看。")
    footer(slide, n, TOTAL)

    risks = [
        ("量产与良率", "A0/B0 实验室成功之后，是 CoWoS、HBM4、铜缆背板、OCS 的供应链和良率。SA 认为这会是高销量 ASIC，但节奏仍在 2027。"),
        ("可编程性债", "OoO + 预取把性能上限抬高，也把「写到 roofline」的难度留给软件。现在靠 Codex 和顶尖内核组，换模型 / 换并行度时能否稳定复现？"),
        ("Agent 负载", "8k1k 能打，不代表百万上下文、多轮 cache、工具调用风暴也能打。这正是 AgentX 要压的东西。"),
        ("同代窗口", "Rubin 已在出货，Jalapeño 大规模要到 2027。窗口期内 Nvidia 软件栈会继续涨。需要的是持续迭代，不是一次发布。"),
        ("组织依赖", "Ho 说内部算力都不够，不会外卖。好处是聚焦，坏处是没有外部客户帮你找软件 bug、也没有第二收入。"),
        ("训练仍外购", "推理自研成功，不降低训练侧对 Nvidia 的依赖。OpenAI 的算力结构会变成「训练买、推理掺着做」。"),
    ]
    for i, (t, d) in enumerate(risks):
        col, row = i % 3, i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.35) + row * Inches(2.7)
        card(slide, x, y, Inches(4.0), Inches(2.5))
        add_textbox(slide, x + Inches(0.22), y + Inches(0.22), Inches(3.55), Inches(0.4),
                    t, 16, True, AMBER)
        add_textbox(slide, x + Inches(0.22), y + Inches(0.7), Inches(3.55), Inches(1.55),
                    d, 13, False, WHITE)


def s23_takeaways(prs, n):
    slide = blank(prs)
    title_block(slide, "TAKEAWAYS", "带走这五条", "如果只记一页，记这一页。")
    footer(slide, n, TOTAL)

    items = [
        ("01", "这是 Serving 芯片", "为 ChatGPT / Agent 的在线推理做的，不是训练 GPU，也不外卖。"),
        ("02", "局部性 > 高峰 FLOPS", "KV 和权重留本地，删固定延迟，让小 batch 也贴 roofline。"),
        ("03", "同质机池是有意的", "不做 PD 分离，是因为流量形状会变，拆池会牺牲 KV 局部性和全局利用率。"),
        ("04", "软件速度是真护城河", "9 个月流片、3 个月 bring-up、Codex 写 kernel。硬件软件共设计，比 ISA 更重要。"),
        ("05", "成绩很强，口径要保守", "打赢 Blackwell 可信；对 Rubin 有方向性优势，但量产、Agent 负载和投机解码都还没闭环。"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.28) + i * Inches(1.08)
        card(slide, Inches(0.5), y, Inches(12.35), Inches(0.96))
        round_rect(slide, Inches(0.7), y + Inches(0.22), Inches(0.7), Inches(0.52), GREEN, None, 0.12)
        add_textbox(slide, Inches(0.7), y + Inches(0.22), Inches(0.7), Inches(0.52),
                    num, 14, True, BLACK, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.6), y + Inches(0.12), Inches(10.9), Inches(0.38),
                    title, 16, True, WHITE)
        add_textbox(slide, Inches(1.6), y + Inches(0.50), Inches(10.9), Inches(0.36),
                    desc, 13, False, MUTED)


def s24_appendix(prs, n):
    slide = blank(prs)
    title_block(slide, "APPENDIX", "口径、来源与阅读建议", "本简报仅基于 2026-08 公开信息，用于内部技术学习，不构成供应或投资判断。")
    footer(slide, n, TOTAL)

    card(slide, Inches(0.5), Inches(1.35), Inches(6.05), Inches(5.55))
    add_textbox(slide, Inches(0.75), Inches(1.52), Inches(5.55), Inches(0.35),
                "主要来源", 16, True, GREEN)
    srcs = [
        "OpenAI：Jalapeño 发布稿；Hot Chips 后的 first results",
        "Hot Chips 2026：You Can Just Build Things … Chips（Ho / Narayanaswami / Leary）",
        "SemiAnalysis：Jalapeño vs Blackwell / Rubin，实验室核对 InferenceX",
        "The Register / TechCrunch / DCD：规格、700W TDP、部署节奏",
        "Broadcom：合作与 Tomahawk 网络口径",
    ]
    for i, t in enumerate(srcs):
        add_textbox(slide, Inches(0.75), Inches(2.05) + i * Inches(0.7), Inches(5.55), Inches(0.65),
                    f"{i+1}.  {t}", 13, False, WHITE)

    card(slide, Inches(6.75), Inches(1.35), Inches(6.1), Inches(5.55))
    add_textbox(slide, Inches(7.0), Inches(1.52), Inches(5.6), Inches(0.35),
                "建议对照着读", 16, True, GREEN)
    recs = [
        "和本仓库「超节点规模」对照：2048 卡全局域 ≈ 高频并行组该有多大。",
        "和「PD 分离」对照：什么时候该拆 P/D，什么时候该保 KV 本地。",
        "和 vLLM 学习对照：continuous batching、paged KV、调度，在 Jalapeño 上被假设为硬件友好原语。",
        "数字优先看 tokens/MW 和 tok/s/user，而不是峰值 PFLOPS。",
        "2027 放量前，把「实验室 A0」和「机房里的机群」分开记账。",
    ]
    for i, t in enumerate(recs):
        add_textbox(slide, Inches(7.0), Inches(2.05) + i * Inches(0.7), Inches(5.6), Inches(0.65),
                    f"{i+1}.  {t}", 13, False, WHITE)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    s01_cover(prs)
    builders = [
        s02_agenda, s03_thesis, s04_what_is, s05_why, s06_timeline,
        s07_section_arch, s08_prefill_decode, s09_locality, s10_microarch, s11_no_pdd,
        s12_section_sys, s13_chip_specs, s14_rack, s15_anatomy, s16_network,
        s17_software, s18_ai_design, s19_bench, s20_caveats,
        s21_strategy, s22_risks, s23_takeaways, s24_appendix,
    ]
    for i, fn in enumerate(builders, start=2):
        fn(prs, i)

    # drop unused nsmap warning by touching
    _ = nsmap
    prs.save(OUT)
    print(f"Wrote {OUT}  slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
